from fpdf import FPDF
from fpdf.enums import XPos, YPos
import os
from datetime import datetime
from utils.sections.summary import generate_section1_summary
from utils.sections.introduction import generate_section2_introduction
from utils.sections.data_overview import generate_section3_data_overview
from utils.sections.methodology import generate_section4_methodology
from utils.sections.cross_domain import generate_section6_cross_domain
from utils.sections.recommendation import generate_section7_recommendations
import textwrap
import re
import streamlit as st

# Custom PDF class to manage page numbers from 2nd page onward
class CustomPDF(FPDF):
    def footer(self):
        if self.page_no() > 1:
            self.set_y(-15)
            self.set_font("DejaVu", "I", 10)
            self.cell(0, 10, f"Page {self.page_no()}", 0, 0, "C")

# def safe_multicell(pdf, text, width=180, font_size=12):
#     pdf.set_font("DejaVu", "", font_size)
#     text = text.replace("**", "").replace("###", "").replace("##", "").replace("#", "")
#     text = text.replace("```", "").replace("🚀", "").replace("🧑", "").replace("🤖", "")
#     paragraphs = text.split("\n")
#     for para in paragraphs:
#         para = para.strip()
#         if not para:
#             continue
#         wrapped_lines = textwrap.wrap(para, width=90)
#         for line in wrapped_lines:
#             pdf.cell(0, 10, line, ln=True)
#         pdf.ln(2)
import re

def clean_insight_text(text):
    # Remove markdown tables completely
    text = re.sub(
        r"(\n\|.+?\|\n\|[-| :]+\|\n(?:\|.*?\|\n?)+)",
        "", text, flags=re.MULTILINE
    )
    # Remove unwanted lines like Analytical Insight or ===
    lines = text.split("\n")
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped or "Analytical Insight" in stripped or re.match(r"=+", stripped):
            continue
        cleaned_lines.append(stripped)
    return "\n".join(cleaned_lines)


def safe_multicell(pdf, text, width=180, font_size=12):
    import textwrap
    import re

    pdf.set_font("DejaVu", "", font_size)

    # --- Clean up unwanted content ---
    text = re.sub(r"(?i)^analytical insight:.*$", "", text, flags=re.MULTILINE)  # remove "Analytical Insight: ..." line
    text = re.sub(r"={3,}|-{3,}", "", text)  # remove "===" or "---" separators
    text = re.sub(r"\*\*+", "", text)  # remove bold markdown
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)  # remove markdown headers like ###

    # Optional cleanup
    text = text.replace("🚀", "").replace("🧑", "").replace("🤖", "")
    
    # --- Render cleaned text ---
    paragraphs = text.split("\n")
    for para in paragraphs:
        para = para.strip()
        if para:
            wrapped_lines = textwrap.wrap(para, width=90)
            for line in wrapped_lines:
                pdf.cell(0, 10, line, ln=True)
            pdf.ln(2)
import re

def extract_markdown_tables(text):
    pattern = r"((\|.+?\|)+\n(\|[-:]+?\|)+\n((\|.*?\|)+\n?)+)"
    return re.findall(pattern, text)

def render_markdown_table(pdf, table_text):
    lines = [line.strip() for line in table_text.strip().split('\n') if line.strip()]
    if not lines or len(lines) < 2:
        return

    headers = [cell.strip() for cell in lines[0].split('|') if cell.strip()]
    rows = [
        [cell.strip() for cell in line.split('|') if cell.strip()]
        for line in lines[2:]
    ]

    col_width = pdf.w / (len(headers) + 1)  # Adjust width

    # Render header
    pdf.set_font("DejaVu", "B", 11)
    for header in headers:
        pdf.cell(col_width, 10, header, border=1)
    pdf.ln()

    # Render rows
    pdf.set_font("DejaVu", "", 11)
    for row in rows:
        for cell in row:
            pdf.cell(col_width, 10, cell, border=1)
        pdf.ln()

    pdf.ln(5)
# def draw_markdown_table(pdf, markdown_text):
#     lines = markdown_text.strip().split("\n")
#     table_lines = [line for line in lines if "|" in line and "---" not in line]

#     if not table_lines or len(table_lines) < 2:
#         return False  # No valid table detected

#     # Extract headers and rows properly
#     headers = [cell.strip() for cell in table_lines[0].split('|') if cell.strip()]
#     rows = [
#         [cell.strip() for cell in line.split('|') if cell.strip()]
#         for line in table_lines[1:]
#     ]

#     col_width = 180 // len(headers)
#     pdf.set_font("DejaVu", "B", 11)

#     # Render header
#     for header in headers:
#         pdf.cell(col_width, 10, header, border=1)
#     pdf.ln()

#     # Render rows
#     pdf.set_font("DejaVu", "", 11)
#     for row in rows:
#         for cell in row:
#             pdf.cell(col_width, 10, cell, border=1)
#         pdf.ln()

#     pdf.ln(5)
#     return True


def draw_markdown_table(pdf, markdown_text):
    lines = markdown_text.strip().split("\n")
    table_lines = [line for line in lines if "|" in line and "---" not in line]

    if not table_lines:
        return False  # No table detected

    # Clean and split table rows
    rows = [line.strip().strip("|").split("|") for line in table_lines]
    col_width = 180 // len(rows[0])
    pdf.set_font("DejaVu", "", 11)
    

    for row in rows:
        for cell in row:
            pdf.cell(col_width, 10, cell.strip(), border=1)
        pdf.ln()
    pdf.ln(5)
    return True

def generate_pdf_report(session, filename="summary.pdf", model_source="groq"):
    df = session["df"]
    dataset_name = session.get("name", "Unnamed Dataset")
    insights = session.get("selected_insight_results", [])
    chat_history = session.get("chat_history", [])

    # Set up PDF
    pdf = CustomPDF()
    font_path = os.path.join("assets", "fonts", "DejaVuSans.ttf")
    font_bold_path = os.path.join("assets", "fonts", "DejaVuSans-Bold.ttf")
    font_italic_path = os.path.join("assets", "fonts", "DejaVuSans-Oblique.ttf")

    pdf.add_font("DejaVu", "", font_path)
    pdf.add_font("DejaVu", "B", font_bold_path)
    pdf.add_font("DejaVu", "I", font_italic_path)

    
    pdf.set_font("DejaVu", "", 14)
    pdf.set_auto_page_break(auto=True, margin=15)

    section_pages = []  # to store (section title, page no)

    # === COVER PAGE ===
    pdf.add_page()
    pdf.set_font("DejaVu", "B", 20)
    pdf.cell(0, 20, "Dynamic Dataset Analysis Report", ln=True, align="C")
    pdf.set_font("DejaVu", "", 16)
    pdf.cell(0, 10, "Generated by Dynamic Impact Tool", ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("DejaVu", "", 12)
    pdf.cell(0, 10, f"Date: {datetime.now().strftime('%d %B %Y')}", ln=True, align="C")
    pdf.cell(0, 10, f"Dataset: {dataset_name}", ln=True, align="C")

    # === INDEX PAGE ===
    pdf.add_page()
    section_pages.append(("Index", pdf.page_no()))

    # Add later after collecting page numbers

    # === SECTION 1 ===
    try:
        summary_text = generate_section1_summary(df, model_source)
        pdf.add_page()
        section_pages.append(("Executive Summary", pdf.page_no()))
        pdf.set_font("DejaVu", "B", 16)
        pdf.cell(0, 10, "1. Executive Summary", ln=True)
        pdf.set_font("DejaVu", "", 12)
        safe_multicell(pdf, summary_text)
    except Exception as e:
        st.error(f"PDF export error in Section 1: {e}")

    # === SECTION 2 ===
    pdf.add_page()
    section_pages.append(("Introduction", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "2. Introduction", ln=True)
    safe_multicell(pdf, generate_section2_introduction(df, model_source))

    # === SECTION 3 ===
    pdf.add_page()
    section_pages.append(("Data Overview", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "3. Data Overview", ln=True)
    safe_multicell(pdf, generate_section3_data_overview(df, model_source))

    # === SECTION 4 ===
    pdf.add_page()
    section_pages.append(("Methodology", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "4. Methodology", ln=True)
    safe_multicell(pdf, generate_section4_methodology(df, model_source))

    # === SECTION 5 ===
    pdf.add_page()
    section_pages.append(("Detailed Analysis & Insights", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "5. Detailed Analysis & Insights", ln=True)
    if insights:
        for idx, insight in enumerate(insights, start=1):
            if idx>1:
             pdf.add_page()
            pdf.set_font("DejaVu", "B", 14)
            pdf.multi_cell(0, 10, f"5.{idx} Insight: {insight['question']}")
            pdf.set_font("DejaVu", "", 12)

            result_text = insight["result"]

            table_rendered = draw_markdown_table(pdf, result_text)
       
            if not table_rendered:
                cleaned_text = clean_insight_text(result_text)
                safe_multicell(pdf, cleaned_text)
            else:
                cleaned_text = result_text
            safe_multicell(pdf, cleaned_text)

            image_path = insight.get("image_path")
            if image_path and isinstance(image_path, str) and os.path.exists(image_path):
              pdf.ln(3)
              try:
                  pdf.image(image_path, w=160)
                  pdf.ln(5)
              except Exception as e:
                  st.warning(f"Couldn't embed image for Insight {idx}: {e}")


            pdf.ln(5)
    else:
        safe_multicell(pdf, "No insights were generated for this dataset.")

    # === SECTION 6 ===
    # === SECTION 6 ===
    try:
     cross_domain_text = generate_section6_cross_domain(df, model_source)
     pdf.add_page()
     section_pages.append(("Cross-Domain Insights", pdf.page_no()))
     pdf.set_font("DejaVu", "B", 16)
     pdf.cell(0, 10, "6. Cross-Domain Insights", ln=True)
     pdf.set_font("DejaVu", "", 12)
     safe_multicell(pdf, cross_domain_text)
    except Exception as e:
    # Check for Groq API token limit error and skip if matched
     if "Groq API error" in str(e) and "rate_limit_exceeded" in str(e):
        st.warning("Section 6 (Cross-Domain Insights) skipped due to Groq API token limit.")
        # Skip adding to section_pages (so it's not added to index)
     else:
        # Show error for debugging
        st.error(f"PDF export error in Section 6: {e}")

    # pdf.add_page()
    # section_pages.append(("Cross-Domain Insights", pdf.page_no()))
    # pdf.set_font("DejaVu", "B", 16)
    # pdf.cell(0, 10, "6. Cross-Domain Insights", ln=True)
    # safe_multicell(pdf, generate_section6_cross_domain(df, model_source))

    # === SECTION 7: Recommendations & Actionable Items ===
    recommendations = generate_section7_recommendations(insights, model_source)
    pdf.add_page()
    section_pages.append(("Recommendations & Actionable Items", pdf.page_no()))
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "7. Recommendations & Actionable Items", ln=True)
    pdf.ln(5)

    if recommendations:
     pdf.set_font("DejaVu", "", 12)
     for idx, rec in enumerate(recommendations, start=1):
        text = (
            f"Insight: {rec.get('Insight', 'N/A')}\n"
            f"Recommended Action: {rec.get('Recommended Action', 'N/A')}\n"
            f"Priority: {rec.get('Priority', 'N/A')}\n"
            f"Owner/Team: {rec.get('Owner', 'N/A')}\n"
            f"Timeline: {rec.get('Timeline', 'N/A')}\n"
            "----------------------------------------------------\n"
        )
        safe_multicell(pdf, text)
    else:
     pdf.set_font("DejaVu", "", 12)
     pdf.multi_cell(0, 10, "No recommendations were generated.")


    # === INDEX PAGE REWRITE ===
    pdf.page = 2  # Set to index page (2nd page)
    pdf.set_xy(10, 20)
    pdf.set_font("DejaVu", "B", 16)
    pdf.cell(0, 10, "Index", ln=True, align="C")
    pdf.ln(5)
    pdf.set_font("DejaVu", "B", 12)
    pdf.cell(20, 10, "S.No", border=1)
    pdf.cell(110, 10, "Section Title", border=1)
    pdf.cell(30, 10, "Page No.", border=1, ln=True)
    pdf.set_font("DejaVu", "", 12)
    for i, (title, page) in enumerate(section_pages[1:], start=1):  # skip index
        pdf.cell(20, 10, str(i), border=1)
        pdf.cell(110, 10, title, border=1)
        pdf.cell(30, 10, str(page), border=1, ln=True)

    # === Export ===
    os.makedirs("exports", exist_ok=True)
    output_path = os.path.join("exports", filename)
    pdf.output(output_path)
    return output_path









# ==================== PPTX EXPORTER ==================== #
from pptx import Presentation
def export_to_pptx(session, filename="summary.pptx"):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"Dataset: {session.get('name', 'Unnamed Dataset')}"
    insight_list = session.get("insights") or session.get("selected_insight_results", [])
    for insight in insight_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"🔍 {insight['question']}"
        slide.shapes.placeholders[1].text = insight['result']

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Chat History"
    chat_text = ""
    for chat in session.get('chat_history', []):
        chat_text += f"🧑 {chat['user']}\n🤖 {chat['assistant'].get('response', '')}\n\n"

    slide.shapes.placeholders[1].text = chat_text

    output_path = os.path.join("exports", filename)
    prs.save(output_path)
    return output_path



